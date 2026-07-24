"""The PPTX parse path — a slide deck serialised to one Markdown page.

A ``.pptx`` is a *flow* format: it has no page grid the way a PDF does, so there
is nothing to measure and route per page. The whole deck therefore collapses onto
a single :class:`~ragsage.models.Page` (number 1, ``layout=None``) whose ``text``
is the deck rendered as Markdown — the same shape the plain-text path emits, so
the classifier's text-presence fallback routes it to ``TEXT`` and the shared
two-pass chunker does the rest.

The serialisation is chosen for that chunker: each slide's title becomes a ``##``
heading (so :func:`~ragsage.parsing.chunking.chunk_markdown` splits on slide
boundaries and carries the slide title into ``metadata["headings"]``), body text
becomes indent-preserving Markdown bullets, tables become GitHub pipe tables (kept
whole by the chunker), and speaker notes are appended under a ``### Notes``
sub-heading so nothing a presenter wrote is dropped from the corpus.

``python-pptx`` is imported **lazily inside** :func:`parse_pptx`, never at module
top level, so importing :mod:`ragsage.parsing` (which the composition root does)
never drags ``pptx``/``lxml``/``pillow`` into the import graph — keeping the
CPU-compatibility dependency guard green.
"""

from __future__ import annotations

import io
from typing import TYPE_CHECKING, cast

from ragsage.models import Page, RawSource

if TYPE_CHECKING:
    from pptx.shapes.base import BaseShape
    from pptx.slide import Slide
    from pptx.table import Table
    from pptx.text.text import TextFrame

# A flow format has no page grid: the whole deck is page 1, layout-free.
_FLOW_FORMAT_PAGE = 1

# One Markdown list indent level is two spaces (the chunker reads standard
# Markdown), so a paragraph nested N levels deep in a text frame is indented 2*N.
_INDENT = "  "


def parse_pptx(source: RawSource, content: bytes) -> list[Page]:
    """Serialise a PPTX deck's bytes to a single Markdown :class:`Page`.

    ``content`` is the raw ``.pptx`` container (``source`` carries only the label);
    it is wrapped in :class:`io.BytesIO` so ``python-pptx`` can open it without a
    file on disk. Slides are walked in order; within each slide the title, body
    text frames, and tables are emitted as Markdown, followed by any speaker notes.
    """
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(content))
    blocks: list[str] = []
    for slide in presentation.slides:
        blocks.extend(_slide_blocks(slide))
    text = "\n\n".join(blocks).strip() + "\n" if blocks else ""
    return [Page(number=_FLOW_FORMAT_PAGE, text=text, layout=None)]


def _slide_blocks(slide: Slide) -> list[str]:
    """The ordered Markdown blocks for one slide: heading, body, tables, notes."""
    title_shape = slide.shapes.title
    title = _clean(title_shape.text) if title_shape is not None else ""
    # Every slide title is an ``##`` heading so the chunker splits on slide
    # boundaries; a titleless slide still gets a heading so its body isn't merged
    # into the previous slide's section.
    blocks: list[str] = [f"## {title}" if title else "##"]

    for shape in slide.shapes:
        if shape is title_shape:
            continue
        if getattr(shape, "has_table", False):
            table = _table_markdown(_shape_table(shape))
            if table:
                blocks.append(table)
        elif getattr(shape, "has_text_frame", False):
            body = _text_frame_markdown(_shape_text_frame(shape))
            if body:
                blocks.append(body)

    notes = _notes_markdown(slide)
    if notes:
        blocks.append(notes)
    return blocks


def _shape_table(shape: BaseShape) -> Table:
    """The ``Table`` behind a shape whose ``has_table`` guard already passed.

    ``python-pptx`` only exposes ``.table`` on the graphic-frame subtype, so mypy
    can't see it on the ``BaseShape`` the ``slide.shapes`` iterator yields; the
    guard is what makes the cast sound.
    """
    return cast("Table", shape.table)  # type: ignore[attr-defined]


def _shape_text_frame(shape: BaseShape) -> TextFrame:
    """The ``TextFrame`` behind a shape whose ``has_text_frame`` guard already passed."""
    return cast("TextFrame", shape.text_frame)  # type: ignore[attr-defined]


def _text_frame_markdown(text_frame: TextFrame) -> str:
    """A body text frame as indent-preserving Markdown bullets.

    Each non-empty paragraph becomes a bullet; ``paragraph.level`` (0-based nesting
    from the PPTX outline) sets the Markdown list indent so nested bullets survive.
    """
    lines: list[str] = []
    for paragraph in text_frame.paragraphs:
        line = _clean(paragraph.text)
        if not line:
            continue
        level = max(getattr(paragraph, "level", 0) or 0, 0)
        lines.append(f"{_INDENT * level}- {line}")
    return "\n".join(lines)


def _table_markdown(table: Table) -> str:
    """A PPTX table as a GitHub pipe table (header row, delimiter, data rows)."""
    rows = [[_cell_text(cell.text) for cell in row.cells] for row in table.rows]
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    header = _pipe_row(rows[0], width)
    delimiter = "| " + " | ".join(["---"] * width) + " |"
    body = [_pipe_row(row, width) for row in rows[1:]]
    return "\n".join([header, delimiter, *body])


def _notes_markdown(slide: Slide) -> str:
    """Speaker notes under a ``### Notes`` sub-heading, or ``""`` if there are none."""
    if not getattr(slide, "has_notes_slide", False):
        return ""
    notes_frame = slide.notes_slide.notes_text_frame
    if notes_frame is None:
        return ""
    notes = _clean(notes_frame.text)
    if not notes:
        return ""
    return f"### Notes\n\n{notes}"


def _pipe_row(cells: list[str], width: int) -> str:
    """Render one table row as ``| a | b |``, padding short rows to ``width`` cells."""
    padded = cells + [""] * (width - len(cells))
    return "| " + " | ".join(padded) + " |"


def _cell_text(text: str) -> str:
    """A table cell flattened to a single line so it can't break the pipe grid.

    A literal ``|`` is escaped so a pipe inside cell text can't be read as a
    column separator, matching how the other format paths serialise cells.
    """
    return " ".join(_clean(text).split()).replace("|", "\\|")


def _clean(text: str) -> str:
    """Trim a run of PPTX text, normalising the vertical-tab soft line breaks it uses."""
    return text.replace("\v", " ").strip()
