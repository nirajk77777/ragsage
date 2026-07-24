"""Coverage for the PPTX parse path (:mod:`ragsage.parsing.pptx`).

The deck is built programmatically with ``python-pptx`` itself — a titled slide
with body bullets, a table, and speaker notes — so the fixture is exactly the
shape the acceptance criteria name, with no binary blob checked in. Two layers are
exercised: :func:`~ragsage.parsing.pptx.parse_pptx` directly (single layout-free
page; title as ``##``; nested bullets; a pipe table; notes present) and the whole
backend (``parse`` then ``chunk``) to prove the slide title reaches
``metadata["headings"]`` and the table survives a real chunking whole.

``python-pptx`` is imported lazily inside the parse path, so a separate guard test
here asserts that importing :mod:`ragsage.parsing` never drags ``pptx`` in.
"""

from __future__ import annotations

import importlib
import io
import sys

from pptx import Presentation
from pptx.util import Inches

from ragsage.models import Page, RawSource
from ragsage.parsing import HeuristicBackend
from ragsage.parsing.pptx import parse_pptx

_PPTX_MEDIA_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
_SLIDE_TITLE = "Quarterly Results"


def _build_deck() -> bytes:
    """A one-slide deck: title, nested body bullets, a table, and speaker notes."""
    presentation = Presentation()
    # A "Title and Content" layout gives us a title placeholder plus a body frame.
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = _SLIDE_TITLE

    body = slide.placeholders[1].text_frame
    body.text = "Revenue rose across every region"
    nested = body.add_paragraph()
    nested.text = "North America led enterprise growth"
    nested.level = 1

    # A 3x2 table (header row + two data rows).
    table = slide.shapes.add_table(3, 2, Inches(1), Inches(3), Inches(6), Inches(2)).table
    for col, label in enumerate(("Item", "Qty")):
        table.cell(0, col).text = label
    table.cell(1, 0).text = "Apple"
    table.cell(1, 1).text = "3"
    table.cell(2, 0).text = "Pear"
    table.cell(2, 1).text = "5"

    slide.notes_slide.notes_text_frame.text = "Remember to thank the finance team."

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _deck_source() -> RawSource:
    return RawSource(name="fixture.pptx", content=_build_deck(), media_type=_PPTX_MEDIA_TYPE)


# --------------------------------------------------------------------------- #
# parse_pptx — a single layout-free Markdown page
# --------------------------------------------------------------------------- #


def test_deck_parses_to_a_single_layout_free_page() -> None:
    pages = parse_pptx(_deck_source(), _build_deck())

    assert len(pages) == 1
    page = pages[0]
    assert isinstance(page, Page)
    assert page.number == 1
    assert page.layout is None  # flow format: classifier falls back to text-presence
    assert page.image is None
    assert page.text.strip()


def test_markdown_captures_title_body_table_and_notes() -> None:
    text = parse_pptx(_deck_source(), _build_deck())[0].text

    # Slide title is a level-2 heading.
    assert f"## {_SLIDE_TITLE}" in text
    # Body text, with the nested bullet indented one Markdown level.
    assert "- Revenue rose across every region" in text
    assert "  - North America led enterprise growth" in text
    # Table rendered as a GitHub pipe grid: header, delimiter, data rows.
    assert "| Item | Qty |" in text
    assert "| --- | --- |" in text
    assert "| Apple | 3 |" in text
    assert "| Pear | 5 |" in text
    # Speaker notes captured under their own sub-heading.
    assert "### Notes" in text
    assert "Remember to thank the finance team." in text


def test_titleless_slide_still_emits_a_heading() -> None:
    # A slide with no title must not merge its body into a neighbour's section.
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # blank layout
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    box.text_frame.text = "Orphan body text"
    buffer = io.BytesIO()
    presentation.save(buffer)
    content = buffer.getvalue()

    text = parse_pptx(RawSource(name="x.pptx", content=content), content)[0].text
    assert text.startswith("##")
    assert "Orphan body text" in text


def test_empty_deck_yields_an_empty_page() -> None:
    presentation = Presentation()
    buffer = io.BytesIO()
    presentation.save(buffer)
    content = buffer.getvalue()

    pages = parse_pptx(RawSource(name="empty.pptx", content=content), content)
    assert len(pages) == 1
    assert pages[0].text == ""
    assert pages[0].layout is None


# --------------------------------------------------------------------------- #
# End-to-end through the backend: headings metadata + tables kept whole
# --------------------------------------------------------------------------- #


def test_backend_routes_pptx_and_carries_slide_title_into_headings() -> None:
    backend = HeuristicBackend()
    parsed = backend.parse(_deck_source())

    # A PPTX is a flow format: it collapses to a single layout-free page.
    assert len(parsed.pages) == 1
    assert parsed.pages[0].layout is None

    chunks = backend.chunk(parsed.document, parsed.pages, size=512, overlap=64)
    assert chunks
    # The slide title rides onto its chunks as heading context.
    assert any(_SLIDE_TITLE in chunk.metadata.get("headings", ()) for chunk in chunks)


def test_backend_keeps_the_table_whole_after_chunking() -> None:
    backend = HeuristicBackend()
    parsed = backend.parse(_deck_source())
    chunks = backend.chunk(parsed.document, parsed.pages, size=512, overlap=64)

    table_lines = ["| Item | Qty |", "| Apple | 3 |", "| Pear | 5 |"]
    holders = [c for c in chunks if all(line in c.text for line in table_lines)]
    assert len(holders) == 1  # the whole table survives inside one chunk


# --------------------------------------------------------------------------- #
# Lazy import guard: importing the parsing package must not drag pptx in
# --------------------------------------------------------------------------- #


def test_importing_parsing_does_not_import_pptx() -> None:
    for name in list(sys.modules):
        if name == "pptx" or name.startswith("pptx."):
            del sys.modules[name]
    for name in list(sys.modules):
        if name == "ragsage.parsing" or name.startswith("ragsage.parsing"):
            del sys.modules[name]

    importlib.import_module("ragsage.parsing")
    assert "pptx" not in sys.modules
